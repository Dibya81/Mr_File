from pptx import Presentation


class PowerPointExtractor:
    def extract(self, file_path: str) -> dict:
        try:
            prs = Presentation(file_path)

            slide_count = len(prs.slides)
            slide_titles = []
            total_text_length = 0
            image_count = 0

            for slide in prs.slides:
                title = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        total_text_length += len(text)
                        if not title and text.strip():
                            title = text.strip()
                    if shape.shape_type == 13:  # Picture
                        image_count += 1
                if title:
                    slide_titles.append(title)

            word_count = total_text_length // 5 if total_text_length > 0 else 0

            core_props = prs.core_properties
            result = {
                "slide_count": slide_count,
                "slide_titles": slide_titles,
                "total_text_length": total_text_length,
                "word_count": word_count,
                "image_count": image_count,
                "title": core_props.title if core_props else "",
                "author": core_props.author if core_props else "",
            }
            return result
        except Exception as e:
            return {"error": str(e), "slide_count": 0, "word_count": 0}


powerpoint_extractor = PowerPointExtractor()
